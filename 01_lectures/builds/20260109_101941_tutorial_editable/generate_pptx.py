#!/usr/bin/env python3
"""
編集可能なPowerPointを生成するスクリプト
白背景、キウイ丸フォント、シンプルデザイン
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# 設定
FONT_NAME = "Kiwi Maru"  # キウイ丸
FONT_NAME_FALLBACK = "Hiragino Sans"  # フォールバック
TITLE_SIZE = Pt(40)
SUBTITLE_SIZE = Pt(24)
BODY_SIZE = Pt(20)
SMALL_SIZE = Pt(16)

# 色
BLUE = RGBColor(21, 101, 192)  # #1565c0
ORANGE = RGBColor(255, 152, 0)  # #ff9800
GREEN = RGBColor(76, 175, 80)  # #4caf50
DARK = RGBColor(51, 51, 51)  # #333333
WHITE = RGBColor(255, 255, 255)

def set_font(run, size=BODY_SIZE, bold=False, color=DARK):
    """フォント設定"""
    run.font.name = FONT_NAME
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color

def add_title_slide(prs, title, subtitle=""):
    """タイトルスライド追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト

    # タイトル
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    set_font(run, TITLE_SIZE, bold=True, color=BLUE)

    # サブタイトル
    if subtitle:
        top = Inches(4.2)
        txBox2 = slide.shapes.add_textbox(left, top, width, Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        set_font(run2, SUBTITLE_SIZE, color=DARK)

    return slide

def add_content_slide(prs, title, bullets, image_path=None):
    """コンテンツスライド追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト

    # タイトル
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, Pt(32), bold=True, color=BLUE)

    # 区切り線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.1), Inches(9), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    # 本文エリア
    content_top = Inches(1.4)
    content_width = Inches(5) if image_path else Inches(9)

    if bullets:
        txBox2 = slide.shapes.add_textbox(left, content_top, content_width, Inches(5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()

            # 箇条書きマーカー
            run = p.add_run()
            run.text = "● " + bullet
            set_font(run, BODY_SIZE, color=DARK)
            p.space_after = Pt(12)

    # 画像
    if image_path and Path(image_path).exists():
        img_left = Inches(5.5)
        img_top = Inches(1.5)
        img_width = Inches(4)
        slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width)

    return slide

def add_section_slide(prs, section_title, section_subtitle=""):
    """セクションタイトルスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景色（薄い青）
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(227, 242, 253)  # #e3f2fd
    background.line.fill.background()

    # セクションタイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = section_title
    set_font(run, TITLE_SIZE, bold=True, color=BLUE)

    if section_subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = section_subtitle
        set_font(run2, SUBTITLE_SIZE, color=DARK)

    return slide

def add_code_slide(prs, title, code_text, explanation=""):
    """コードスライド追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = title
    set_font(run, Pt(32), bold=True, color=BLUE)

    # コードボックス
    code_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.3), Inches(9), Inches(3.5)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(245, 245, 245)  # 薄いグレー
    code_box.line.color.rgb = RGBColor(200, 200, 200)

    # コードテキスト
    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(3.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    run = p.add_run()
    run.text = code_text
    run.font.name = "Menlo"
    run.font.size = Pt(14)
    run.font.color.rgb = DARK

    # 説明
    if explanation:
        txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = explanation
        set_font(run3, BODY_SIZE, color=DARK)

    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """画像スライド追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = title
    set_font(run, Pt(32), bold=True, color=BLUE)

    # 画像
    if Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), Inches(1.5), Inches(1.3), width=Inches(7))

    # キャプション
    if caption:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
        tf2 = txBox2.text_frame
        p = tf2.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run2 = p.add_run()
        run2.text = caption
        set_font(run2, SMALL_SIZE, color=DARK)

    return slide


def main():
    """メイン処理"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    diagrams = Path("diagrams")

    # ===== タイトル =====
    add_title_slide(prs,
        "AI × Game Development Tutorial",
        "はじめてでも分かる「ゲームが動く仕組み」と「AIが考える仕組み」"
    )

    # ===== 導入 =====
    add_content_slide(prs, "この資料について", [
        "「ゲームを作ったことがない人」向け",
        "「AIやプログラミングが初めての人」向け",
        "分からない言葉は必ず説明します",
        "順番に分かるようになります"
    ])

    add_content_slide(prs, "この教材の本当の目的", [
        "ゲームを完成させることではありません",
        "",
        "ゲームやAIが",
        "「なぜその動きをしたのか」を",
        "自分の言葉で説明できるようになること"
    ])

    # 全体像
    add_image_slide(prs, "全体像：9つのStep",
        diagrams / "steps_overview.png",
        "基礎 → DSL → AI の順に学びます"
    )

    # ===== Step 00 =====
    add_section_slide(prs, "Step 00", "Hello World - 環境確認")

    add_content_slide(prs, "Step 00: 目的", [
        "Pythonが動くことを確認",
        "ゲームを起動してみる",
        "w/a/s/d で動かしてみる"
    ])

    add_code_slide(prs, "Step 00: やってみよう",
        "cd 02_tutorial\npython3 run.py\n\n1. 「New Game」を選ぶ\n2. 「Step 00」を選ぶ\n3. SAVE_A に保存\n4. w を押す",
        "@ が上に動いたら成功です！"
    )

    # ===== Step 01 =====
    add_section_slide(prs, "Step 01", "Game Loop - ゲームループ")

    add_image_slide(prs, "ゲームループとは？",
        diagrams / "game_loop.png",
        "入力 → 更新 → 描画 を繰り返す"
    )

    add_content_slide(prs, "Step 01: カウンターゲーム", [
        "up / + : 数字が増える",
        "down / - : 数字が減る",
        "reset : 0に戻る",
        "",
        "これがゲームの本質！"
    ])

    add_content_slide(prs, "なぜカウンター？", [
        "RPGでHPが減るのも",
        "スコアが増えるのも",
        "レベルが上がるのも",
        "",
        "全部「状態の数値を変える」だけ！"
    ])

    # ===== Step 02 =====
    add_section_slide(prs, "Step 02", "State - 状態管理")

    add_content_slide(prs, "「状態」って何？", [
        "ゲームの「今どうなってるか」を表すデータ",
        "",
        "例：",
        "  player: x=10, y=5, hp=100",
        "  turn: 0",
        "  score: 0"
    ])

    add_image_slide(prs, "イミュータブル（不変）",
        diagrams / "state_immutable.png",
        "古い状態を変更せず、新しい状態を作る"
    )

    add_content_slide(prs, "イミュータブルのメリット", [
        "何が起きたか追跡しやすい",
        "バグの原因を見つけやすい",
        "「Undo」も簡単に作れる"
    ])

    # ===== Step 03 =====
    add_section_slide(prs, "Step 03", "I/O - 入出力の分離")

    add_image_slide(prs, "入出力の分離",
        diagrams / "io_separation.png",
        "ゲームロジックを「純粋」にする"
    )

    add_content_slide(prs, "純粋関数とは？", [
        "同じ入力なら必ず同じ出力",
        "副作用がない（外部状態を変更しない）",
        "",
        "例：add(2, 3) は常に 5 を返す"
    ])

    # ===== Step 04 =====
    add_section_slide(prs, "Step 04", "Renderer - レンダラー")

    add_content_slide(prs, "レンダラーとは？", [
        "状態を「見える形」に変換する関数",
        "",
        "入力：ゲームの状態（state）",
        "出力：描画結果の文字列",
        "副作用なし：print しない！"
    ])

    add_code_slide(prs, "出力例",
        "######################\n#..........@.........#\n#....................#\n#........E...........#\n#....................#\n######################",
        "@ = プレイヤー、E = 敵"
    )

    # ===== DSLセクション =====
    add_section_slide(prs, "DSL編", "Step 05-07: ゲーム言語を作る")

    add_image_slide(prs, "DSLの処理フロー",
        diagrams / "dsl_flow.png",
        "文字列 → トークン → AST → 実行"
    )

    # ===== Step 05 =====
    add_section_slide(prs, "Step 05", "Lexer - 字句解析")

    add_content_slide(prs, "Lexerとは？", [
        "文字列を「単語（トークン）」に分解",
        "",
        "例：move player 5 3",
        "→ [move] [player] [5] [3]",
        "",
        "それぞれに「種類」をつける"
    ])

    # ===== Step 06 =====
    add_section_slide(prs, "Step 06", "Parser - 構文解析")

    add_content_slide(prs, "Parserとは？", [
        "トークンを「構文木（AST）」に変換",
        "",
        "例：move player 5 3",
        "→ MoveCommand(",
        "      target=\"player\",",
        "      x=5, y=3)"
    ])

    # ===== Step 07 =====
    add_section_slide(prs, "Step 07", "Interpreter - インタプリタ")

    add_content_slide(prs, "Interpreterとは？", [
        "ASTを「実行」して状態を変更",
        "",
        "使えるコマンド：",
        "  move player <x> <y>",
        "  spawn enemy <x> <y>",
        "  destroy <target>",
        "  wait（AIターン）"
    ])

    # ===== Step 08 =====
    add_section_slide(prs, "Step 08", "AI - A*経路探索")

    add_content_slide(prs, "A*アルゴリズムとは？", [
        "「最短経路」を見つけるアルゴリズム",
        "",
        "敵がプレイヤーを追いかける仕組み"
    ])

    add_image_slide(prs, "A*の処理フロー",
        diagrams / "astar.png",
        "ゴールに近い場所を優先して探索"
    )

    add_content_slide(prs, "AIは「考えて」いるのか？", [
        "実際はとてもシンプル：",
        "",
        "1. 現在位置とゴールを確認",
        "2. 行ける場所をリストアップ",
        "3. ゴールに近い場所を優先",
        "4. 繰り返す",
        "",
        "「ルール通りに計算している」だけ！"
    ])

    # ===== まとめ =====
    add_section_slide(prs, "まとめ", "")

    add_content_slide(prs, "この教材で学んだこと", [
        "1. ゲームループ：入力→更新→描画",
        "2. 状態管理：イミュータブル",
        "3. DSL：Lexer → Parser → Interpreter",
        "4. AI：ルール通りに計算"
    ])

    add_content_slide(prs, "最後の質問", [
        "",
        "「なぜその動きをしたのか」",
        "説明できるようになりましたか？",
        "",
        "できるようになったら、あなたはもう",
        "ゲームとAIの「仕組み」を理解した人です！"
    ])

    add_title_slide(prs, "Thank You!", "AI × Game Development Tutorial")

    # 保存
    output_path = Path("tutorial_slides.pptx")
    prs.save(output_path)
    print(f"Generated: {output_path}")
    print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()

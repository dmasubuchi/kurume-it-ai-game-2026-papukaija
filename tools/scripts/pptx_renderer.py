#!/usr/bin/env python3
"""
PowerPoint Renderer
スライドマスターを活用してJSONからPowerPointファイルを生成
"""

import json
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


class PowerPointRenderer:
    """PowerPoint生成クラス（スライドマスター対応）"""

    # スライドレイアウトのインデックス
    LAYOUT_TITLE = 0          # タイトルスライド
    LAYOUT_TITLE_CONTENT = 1  # タイトルと本文
    LAYOUT_SECTION = 2        # セクションヘッダー
    LAYOUT_TWO_CONTENT = 3    # 2つのコンテンツ
    LAYOUT_COMPARISON = 4     # 比較
    LAYOUT_TITLE_ONLY = 5     # タイトルのみ
    LAYOUT_BLANK = 6          # 白紙

    def __init__(self, theme: dict, project_root: str):
        self.theme = theme
        self.project_root = Path(project_root)
        self.processing_log = []

        # スライドサイズ
        slide_config = theme.get("slide", {})
        self.slide_width = Inches(slide_config.get("width_inches", 13.333))
        self.slide_height = Inches(slide_config.get("height_inches", 7.5))

        # 色設定
        colors = theme.get("colors", {})
        self.bg_color = self._parse_color(colors.get("background", "#FFFFFF"))
        self.bg_dark = self._parse_color(colors.get("background_dark", "#113178"))
        self.primary = self._parse_color(colors.get("primary", "#113178"))
        self.text_color = self._parse_color(colors.get("text", "#333333"))
        self.text_light = self._parse_color(colors.get("text_light", "#FFFFFF"))
        self.text_secondary = self._parse_color(colors.get("text_secondary", "#666666"))

        # フォント設定
        fonts = theme.get("fonts", {})
        self.heading_font = fonts.get("heading", {}).get("name", "游ゴシック")
        self.body_font = fonts.get("body", {}).get("name", "游ゴシック")
        self.heading_size = Pt(fonts.get("heading", {}).get("size", 32))
        self.body_size = Pt(fonts.get("body", {}).get("size", 20))
        self.title_slide_size = Pt(fonts.get("title_slide", {}).get("size", 44))

    def _parse_color(self, color_str: str) -> RGBColor:
        """カラー文字列をRGBColorに変換"""
        color_str = color_str.lstrip("#")
        return RGBColor(
            int(color_str[0:2], 16),
            int(color_str[2:4], 16),
            int(color_str[4:6], 16)
        )

    def _set_japanese_font(self, run, font_name: str):
        """日本語フォントを設定"""
        run.font.name = font_name
        r = run._r
        rPr = r.get_or_add_rPr()
        latin = rPr.get_or_add_latin()
        latin.set("typeface", font_name)
        from lxml import etree
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font_name)

    def _setup_slide_master(self, prs):
        """スライドマスターのスタイルを設定"""
        slide_master = prs.slide_master

        # マスターの各レイアウトに対してフォントと色を設定
        for layout in prs.slide_layouts:
            for shape in layout.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            self._set_japanese_font(run, self.heading_font)

    def render(self, data: dict, output_path: str):
        """プレゼンテーションを生成"""
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height

        slides_data = data.get("slides", [])
        assets = data.get("assets", {})

        print(f"生成開始: {len(slides_data)}枚のスライド")

        for i, slide_data in enumerate(slides_data, 1):
            layout = slide_data.get("layout", "content")

            if layout == "title":
                self._create_title_slide(prs, slide_data, assets)
            elif layout == "section":
                self._create_section_slide(prs, slide_data, assets)
            elif layout == "two_column":
                self._create_two_column_slide(prs, slide_data, assets)
            elif layout == "visual":
                self._create_visual_slide(prs, slide_data, assets)
            else:
                self._create_content_slide(prs, slide_data, assets)

            if i % 20 == 0:
                print(f"  {i}/{len(slides_data)} 完了...")

        prs.save(output_path)
        print(f"\n生成完了: {output_path}")
        print(f"総スライド数: {len(prs.slides)}")

        return {
            "output_path": output_path,
            "slide_count": len(prs.slides),
            "processing_log": self.processing_log
        }

    def _create_title_slide(self, prs, slide_data: dict, assets: dict):
        """タイトルスライドを作成（レイアウト0: タイトルスライド）"""
        slide_layout = prs.slide_layouts[self.LAYOUT_TITLE]
        slide = prs.slides.add_slide(slide_layout)

        # 背景色を設定
        self._set_slide_background(slide, self.bg_dark)

        # タイトルプレースホルダーを使用（明示的にフォント設定）
        title_shape = slide.shapes.title
        if title_shape:
            for elem in slide_data.get("elements", []):
                if elem["type"] == "title":
                    self._set_text_with_font(
                        title_shape, elem["content"],
                        self.heading_font, self.title_slide_size, self.text_light, bold=True
                    )
                    break

        # サブタイトルプレースホルダーを使用
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # サブタイトル
                for elem in slide_data.get("elements", []):
                    if elem["type"] == "subtitle" or elem["type"] == "text":
                        self._set_text_with_font(
                            shape, elem["content"],
                            self.body_font, Pt(24), self.text_secondary, align=PP_ALIGN.CENTER
                        )
                        break

        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    def _create_section_slide(self, prs, slide_data: dict, assets: dict):
        """セクション区切りスライドを作成（レイアウト2: セクションヘッダー）"""
        slide_layout = prs.slide_layouts[self.LAYOUT_SECTION]
        slide = prs.slides.add_slide(slide_layout)

        # 背景色を設定
        self._set_slide_background(slide, self.bg_dark)

        # タイトルプレースホルダーを使用（明示的にフォント設定）
        title_shape = slide.shapes.title
        if title_shape:
            for elem in slide_data.get("elements", []):
                if elem["type"] == "title":
                    self._set_text_with_font(
                        title_shape, elem["content"],
                        self.heading_font, Pt(40), self.text_light, bold=True
                    )
                    break

        # サブタイトル/テキスト
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                for elem in slide_data.get("elements", []):
                    if elem["type"] == "subtitle" or elem["type"] == "text":
                        self._set_text_with_font(
                            shape, elem["content"],
                            self.body_font, Pt(24), RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER
                        )
                        break

        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    def _create_content_slide(self, prs, slide_data: dict, assets: dict):
        """コンテンツスライドを作成（レイアウト1: タイトルと本文）"""
        slide_layout = prs.slide_layouts[self.LAYOUT_TITLE_CONTENT]
        slide = prs.slides.add_slide(slide_layout)

        # 背景色を設定（白）
        self._set_slide_background(slide, self.bg_color)

        # アクセントライン追加
        self._add_accent_line(slide)

        # タイトルプレースホルダーを使用（明示的にフォント設定）
        title_shape = slide.shapes.title
        if title_shape:
            for elem in slide_data.get("elements", []):
                if elem["type"] == "title":
                    self._set_text_with_font(
                        title_shape, elem["content"],
                        self.heading_font, self.heading_size, self.primary, bold=True
                    )
                    break

        # 本文プレースホルダーを使用
        body_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # 本文プレースホルダー
                body_shape = shape
                break

        if body_shape:
            tf = body_shape.text_frame
            tf.clear()
            first_para = True

            for elem in slide_data.get("elements", []):
                if elem["type"] == "bullet_list":
                    for item in elem["content"]:
                        if first_para:
                            p = tf.paragraphs[0]
                            first_para = False
                        else:
                            p = tf.add_paragraph()
                        p.level = 0
                        # 明示的にrunを作成してフォントを設定
                        run = p.add_run()
                        run.text = item
                        self._set_japanese_font(run, self.body_font)
                        run.font.size = self.body_size
                        run.font.color.rgb = self.text_color
                        p.space_before = Pt(8)
                        p.space_after = Pt(4)

                elif elem["type"] == "text" and elem["type"] != "title":
                    if first_para:
                        p = tf.paragraphs[0]
                        first_para = False
                    else:
                        p = tf.add_paragraph()
                    p.level = 0
                    # 明示的にrunを作成してフォントを設定
                    run = p.add_run()
                    run.text = elem["content"]
                    self._set_japanese_font(run, self.body_font)
                    run.font.size = self.body_size
                    run.font.color.rgb = self.text_color
                    p.space_before = Pt(8)
                    p.space_after = Pt(4)

                elif elem["type"] == "diagram" or elem["type"] == "image":
                    # 画像は本文の下に追加
                    self._add_image_to_slide(slide, elem, assets)

        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    def _create_two_column_slide(self, prs, slide_data: dict, assets: dict):
        """2カラムスライドを作成（レイアウト3: 2つのコンテンツ）"""
        slide_layout = prs.slide_layouts[self.LAYOUT_TWO_CONTENT]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, self.bg_color)
        self._add_accent_line(slide)

        # タイトル（明示的にフォント設定）
        title_shape = slide.shapes.title
        if title_shape:
            for elem in slide_data.get("elements", []):
                if elem["type"] == "title":
                    self._set_text_with_font(
                        title_shape, elem["content"],
                        self.heading_font, self.heading_size, self.primary, bold=True
                    )
                    break

        # 左右のプレースホルダーを探す
        left_shape = None
        right_shape = None
        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx == 1:
                left_shape = shape
            elif idx == 2:
                right_shape = shape

        # 左カラム
        if left_shape:
            for elem in slide_data.get("elements", []):
                if elem["type"] == "left_content":
                    self._fill_content_placeholder(left_shape, elem.get("items", []))
                    break

        # 右カラム
        if right_shape:
            for elem in slide_data.get("elements", []):
                if elem["type"] == "right_content":
                    self._fill_content_placeholder(right_shape, elem.get("items", []))
                    break

        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    def _create_visual_slide(self, prs, slide_data: dict, assets: dict):
        """画像中心スライドを作成（レイアウト5: タイトルのみ + 画像）"""
        slide_layout = prs.slide_layouts[self.LAYOUT_TITLE_ONLY]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, self.bg_color)
        self._add_accent_line(slide)

        # タイトル（明示的にフォント設定）
        title_shape = slide.shapes.title
        if title_shape:
            for elem in slide_data.get("elements", []):
                if elem["type"] == "title":
                    self._set_text_with_font(
                        title_shape, elem["content"],
                        self.heading_font, self.heading_size, self.primary, bold=True
                    )
                    break

        # 画像
        for elem in slide_data.get("elements", []):
            if elem["type"] == "diagram" or elem["type"] == "image":
                self._add_image_to_slide(slide, elem, assets, y_offset=Inches(1.5))
                break

        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    def _set_slide_background(self, slide, color: RGBColor):
        """スライドの背景色を設定"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_accent_line(self, slide):
        """上部アクセントラインを追加"""
        decorations = self.theme.get("decorations", {})
        if not decorations.get("accent_line", {}).get("enabled", True):
            return

        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.slide_width, Inches(0.08)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = self.primary
        accent_line.line.fill.background()

    def _set_text_with_font(self, shape, text: str, font_name: str, font_size, color, bold: bool = False, align=None):
        """テキストを設定し、明示的にフォントを適用"""
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        if align:
            p.alignment = align
        run = p.add_run()
        run.text = text
        self._set_japanese_font(run, font_name)
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color

    def _format_title_placeholder(self, shape, font_size, color):
        """タイトルプレースホルダーのフォーマット（後方互換用）"""
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = font_size
            paragraph.font.bold = True
            paragraph.font.color.rgb = color
            for run in paragraph.runs:
                self._set_japanese_font(run, self.heading_font)
                run.font.size = font_size
                run.font.bold = True
                run.font.color.rgb = color

    def _format_body_placeholder(self, shape, font_size, color):
        """本文プレースホルダーのフォーマット"""
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = font_size
            paragraph.font.color.rgb = color
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                self._set_japanese_font(run, self.body_font)
                run.font.size = font_size
                run.font.color.rgb = color

    def _format_paragraph(self, paragraph, font_size, color):
        """段落のフォーマット"""
        paragraph.font.size = font_size
        paragraph.font.color.rgb = color
        paragraph.space_before = Pt(8)
        paragraph.space_after = Pt(4)
        for run in paragraph.runs:
            self._set_japanese_font(run, self.body_font)
            run.font.size = font_size
            run.font.color.rgb = color

    def _fill_content_placeholder(self, shape, items):
        """コンテンツプレースホルダーを埋める（明示的にフォント設定）"""
        tf = shape.text_frame
        tf.clear()
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.level = 0
            # 明示的にrunを作成してフォントを設定
            run = p.add_run()
            run.text = item
            self._set_japanese_font(run, self.body_font)
            run.font.size = self.body_size
            run.font.color.rgb = self.text_color
            p.space_before = Pt(8)
            p.space_after = Pt(4)

    def _add_image_to_slide(self, slide, elem, assets, y_offset=Inches(2.0)):
        """スライドに画像を追加"""
        asset_ref = elem.get("asset_ref")
        if not asset_ref:
            return

        # 対応するアセットを探す
        diagrams = assets.get("diagrams", [])
        images = assets.get("images", [])

        image_path = None
        for diagram in diagrams:
            if diagram.get("id") == asset_ref:
                image_path = diagram.get("output_path")
                break

        if not image_path:
            for image in images:
                if image.get("id") == asset_ref:
                    image_path = image.get("output_path")
                    break

        if image_path:
            full_path = self.project_root / image_path
            if full_path.exists():
                # 中央配置
                img_width = Inches(8)
                img_left = (self.slide_width - img_width) / 2
                try:
                    slide.shapes.add_picture(
                        str(full_path),
                        img_left,
                        y_offset,
                        width=img_width
                    )
                except Exception as e:
                    print(f"  警告: 画像追加失敗 {asset_ref}: {e}")


def render_presentation(json_path: str, output_path: str, theme: dict, project_root: str):
    """プレゼンテーションをレンダリング"""
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    # presentation キーがある場合はその中身を使用
    if "presentation" in data:
        data = data["presentation"]

    renderer = PowerPointRenderer(theme, project_root)
    return renderer.render(data, output_path)


if __name__ == "__main__":
    import argparse

    arg_parser = argparse.ArgumentParser(description="PowerPoint Renderer")
    arg_parser.add_argument("--input", "-i", required=True, help="入力JSONファイル")
    arg_parser.add_argument("--output", "-o", required=True, help="出力PowerPointファイル")
    arg_parser.add_argument("--theme", "-t", required=True, help="テーマYAMLファイル")
    arg_parser.add_argument("--project", "-p", default=".", help="プロジェクトルート")

    args = arg_parser.parse_args()

    import yaml
    with open(args.theme, "r", encoding="utf-8") as f:
        theme = yaml.safe_load(f)

    render_presentation(args.input, args.output, theme, args.project)
